#!/usr/bin/env python3
"""Build FORTRESS presentation from Alfa Bank Google Slides template (PPTX export)."""

from __future__ import annotations

import copy
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "alfa_template.pptx"
DIAGRAMS = ROOT / "diagrams"
OUTPUT = ROOT / "FORTRESS_защита.pptx"


def set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    shape.text_frame.text = text


def replace_largest_picture(slide, image_path: Path, *, full_bleed: bool = False) -> None:
    if full_bleed:
        # Wide diagrams (сквозной флоу) — на всю ширину слайда, без обрезки
        slide.shapes.add_picture(
            str(image_path),
            Emu(120000),
            Emu(1500000),
            width=Emu(17900000),
            height=Emu(8200000),
        )
        return
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        slide.shapes.add_picture(str(image_path), Emu(200000), Emu(1800000), width=Emu(17000000))
        return
    target = max(pics, key=lambda s: s.width * s.height)
    slide.shapes.add_picture(
        str(image_path),
        target.left,
        target.top,
        width=target.width,
        height=target.height,
    )


def set_table_row(table, row: int, values: list[str]) -> None:
    for col, val in enumerate(values):
        if col < len(table.columns):
            table.cell(row, col).text = val


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def delete_slides(prs: Presentation, indices: list[int]) -> None:
    for idx in sorted(indices, reverse=True):
        delete_slide(prs, idx)


def update_threats_table(slide) -> None:
    """Правим gate-имена в существующей таблице, не урезаем строки."""
    fixes = {
        "G1 баланс+схема+PII": "DATA gate",
        "G4 allow-list форматов + modelscan": "G5 modelscan + G6",
        "G2 gitleaks+pip-audit+bandit+trivy": "G0 gitleaks + G3 pip-audit",
        "G4+Registry SHA-256 + cosign": "G7 SHA-256 + cosign",
        "G7 rate-limit→429 + output-reduction": "G14 rate-limit + DLP",
    }
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        t = shape.table
        for r in range(len(t.rows)):
            for c in range(len(t.columns)):
                val = t.cell(r, c).text
                for old, new in fixes.items():
                    if old in val:
                        t.cell(r, c).text = val.replace(old, new)


def update_analogs_table(slide) -> None:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        t = shape.table
        rows = [
            ["Категория", "Примеры", "Проблема / отличие FORTRESS"],
            ["Платформы", "HiddenLayer, Protect AI, SageMaker", "Дорого / облако / нет bank CI→prod policy"],
            ["Инструменты", "gitleaks, trivy, modelscan, MLflow", "Хороши сами по себе, но НЕ интегрированы"],
            ["FORTRESS", "Self-hosted compose", "Единый pipeline + lineage + audit + Human Approve"],
        ]
        for r in range(min(len(t.rows), len(rows))):
            for c in range(min(len(t.columns), len(rows[r]))):
                t.cell(r, c).text = rows[r][c]


def apply_updates(prs: Presentation) -> None:
    d = DIAGRAMS

    # 1 — Title
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and "ML Sec" in shape.text_frame.text:
            set_text(shape, "FORTRESS\nMLSecOps-платформа для банка")

    # 2 — Problem (short)
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame and "Бизнес-ценность" in shape.text_frame.text:
            set_text(
                shape,
                "Проблема\n\n«ML приносит деньги — пока не ломается»\n\n"
                "Сотни моделей в prod · один сбой = простой скоринга / штраф",
            )

    # 3 — Pipeline + threats (was «ментор»)
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame:
            set_text(shape, "Где возникают угрозы")
    replace_largest_picture(slide, d / "01_trust_zones.png")

    # 4 — Context (4-box problem slide)
    slide = prs.slides[3]
    texts = {
        "Проблема": "Контекст",
        "01": "01",
        "Каждый день": "Сотни ML-моделей\nв prod каждый день",
        "02": "02",
        "ML/DS работает": "Критические активы:\nданные · модели · код",
        "03": "03",
        "У MLSecOps": "MLSecOps тонет\nв рутине проверок",
        "04": "04",
        "Наша платформа": "FORTRESS дополняет\nMLOps — не заменяет ИБ",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            for key, val in texts.items():
                if t.startswith(key) or t == key:
                    set_text(shape, val)
                    break

    # 5 — Assets
    slide = prs.slides[4]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "Основные активы" in t or "Активы под" in t:
            set_text(shape, "Активы под защитой")
        elif "Данные являются" in t:
            set_text(shape, "Данные · PII · 152-ФЗ")
        elif "Постоянные сливы" in t:
            set_text(shape, "Утечки PII\nи секретов")
        elif "основной функционал" in t:
            set_text(shape, "Бизнес завязан\nна ML-моделях")
        elif "стабильная" in t:
            set_text(shape, "Предсказуемость\nи версии")
        elif shape.text_frame.text.strip() == "Данные":
            set_text(shape, "Данные")
        elif "ML модели" in t or shape.text_frame.text.strip() == "ML-модели":
            set_text(shape, "ML-модели")
        elif "Код" in t:
            set_text(shape, "Код · секреты")

    # 6 — Benefits / TTM
    slide = prs.slides[5]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Что улучшает" in t:
            set_text(shape, "Жизненный цикл модели")
        elif "Меньше сбоев" in t:
            set_text(shape, "Меньше инцидентов\nв prod")
        elif "Меньше ручной" in t and "Time" not in t:
            set_text(shape, "Меньше рутины\nMLSecOps")
        elif "Time To Market" in t:
            set_text(shape, "Быстрее\nв Production")
        elif "Автоматический контроль" in t:
            set_text(shape, "Автопроверки\nна каждом этапе")
        elif "человеческого фактора" in t:
            set_text(shape, "Human Approve\nтолько для external")
        elif "цикл жизни" in t:
            set_text(shape, "Полный контур:\nDEV → CI → Prod")
    replace_largest_picture(slide, d / "06_model_lifecycle.png")

    # 7 — Metrics (was 10x slide)
    slide = prs.slides[6]
    mapping = {
        "Что дает": "Метрики платформы",
        "10х": "G0–G15",
        "~0": "3 модели",
        "Гарантия": "Audit chain ✓",
        "Ускорение": "CI-ready",
        "Меньше рутины": "Findings",
        "Случайных": "Human Approve\n(исключение)",
        "продвинутые": "Signed\nattestation",
        "Human Approve": "15 мин vs\n2–5 дней review",
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        for key, val in mapping.items():
            if key in txt:
                set_text(shape, val)
                break

    # 9 — Onboarding
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and "Простой" in shape.text_frame.text:
            set_text(shape, "Онбординг\n\nРегистрация → роль → MLflow SSO\nОдин логин для UI и MLflow")

    # 10 — Jupyter + MLflow
    slide = prs.slides[9]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Jupiter" in t or "Дефолтная" in t:
            set_text(shape, "Jupyter — эксперименты\nMLflow — реестр")
        elif "ML Flow" in t or "MLflow" in t.lower():
            set_text(shape, "В prod — только модель\nиз CI pipeline (signed)")
        elif "не отвлекаться" in t:
            set_text(shape, "Lineage: dataset SHA ↔ git ↔ run_id")
        elif "синхронизация" in t.lower():
            set_text(shape, "Автосинхронизация\nattestation → MLflow tags")

    # 13 — Asset registration
    slide = prs.slides[12]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Автоматическая регистрация" in t:
            set_text(shape, "Все активы через контур FORTRESS\n(MLflow + gates + audit)")
        elif "Разработка" in t:
            set_text(shape, "Эксперименты")
        elif "Подгрузка" in t:
            set_text(shape, "Ручная загрузка")
        elif "безопастности" in t or "безопасности" in t:
            set_text(shape, "Автопроверки")
        elif "Версионирование" in t:
            set_text(shape, "Версии + SHA")
        elif "контур" in t:
            set_text(shape, "Единый audit trail")

    # 14 — MLSecOps pains
    slide = prs.slides[13]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Какая выгода" in t:
            set_text(shape, "Боли MLSecOps → FORTRESS")
        elif "рутинные процессы" in t:
            set_text(shape, "«Где модель?» → единый реестр")
        elif "контроль над безопасностью" in t:
            set_text(shape, "«Прошла проверка?» → теги security.*")
        elif "назначать роли" in t:
            set_text(shape, "«Кто одобрил?» → audit chain")
        elif "HITL" in t:
            set_text(shape, "Рутина автоматизирована\nHuman Approve — исключение")
    replace_largest_picture(slide, d / "10_rbac.png")

    # 16 — Old analogies intro → platforms vs tools visual
    slide = prs.slides[15]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Что уже существует" in t:
            set_text(shape, "Рынок сегодня")
        elif "Энтерпрайз" in t:
            set_text(shape, "Платформы\n(HiddenLayer, Protect AI)")
        elif "Runtime" in t:
            set_text(shape, "Облачный MLOps\n(SageMaker, Vertex)")
        elif "Опен-сорс" in t:
            set_text(shape, "Разрозненные\nинструменты")
        elif "HiddenLayer" in t:
            set_text(shape, "Дорого · vendor lock-in")
        elif "gitleaks" in t:
            set_text(shape, "gitleaks · trivy · modelscan\nне связаны между собой")
        elif "инференс" in t:
            set_text(shape, "Нет lineage · нет «можно в prod»")
    replace_largest_picture(slide, d / "11_analogs_market.png")

    # 17-18 — Team (keep names, light touch)
    for shape in prs.slides[16].shapes:
        if shape.has_text_frame and "Николай" in shape.text_frame.text:
            pass  # keep
    for shape in prs.slides[17].shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Кто чем" in t:
            set_text(shape, "Команда")

    # 19 — Section: threats
    for shape in prs.slides[18].shapes:
        if shape.has_text_frame and "Закрываемые" in shape.text_frame.text:
            set_text(shape, "Модель угроз")
    replace_largest_picture(prs.slides[18], d / "09_threats_mapping.png")

    # 20 — Threats table
    update_threats_table(prs.slides[19])

    # 21 — Tools: таблица шаблона + диаграмма gates из docs §7
    slide = prs.slides[20]
    for shape in slide.shapes:
        if shape.has_text_frame and "Реализованные" in shape.text_frame.text:
            set_text(shape, "Security Gates — слои (docs/architecture.md §7)")
    replace_largest_picture(slide, d / "03_security_gates_layers.png")

    # 22 — Human Approve (was empty header)
    slide = prs.slides[21]
    for shape in slide.shapes:
        if shape.has_text_frame:
            set_text(shape, "Human Approve\n\nТолько для external-моделей и tier HIGH\nCI-модель + signed → deploy без ручного шага")

    # 23-24 — Demo + UI visibility
    for shape in prs.slides[22].shapes:
        t = shape.text_frame.text.strip() if shape.has_text_frame else ""
        if shape.has_text_frame and ("Демо" in t or "Заголовок" in t):
            set_text(shape, "Демо — 3 сценария")
    replace_largest_picture(prs.slides[22], d / "07_demo_scenarios.png")

    for shape in prs.slides[23].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "UI":
            set_text(shape, "Видимость для MLSecOps")
    replace_largest_picture(prs.slides[23], d / "08_visibility_er.png")

    # 25 — Layers summary
    slide = prs.slides[24]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "слои проверок" in t.lower():
            set_text(shape, "5 слоёв: данные · код · модель · образ · runtime")
        elif "Мониторинг" in t:
            set_text(shape, "Streamlit Security Center :8502")

    # 27 — Functionality
    slide = prs.slides[26]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if t == "Функционал":
            set_text(shape, "Функционал FORTRESS")
        elif "Админ" in t:
            set_text(shape, "RBAC / Роли")
        elif "главный специалист" in t:
            set_text(shape, "MLSecOps назначает роли\n(ds · de · mlsecops)")
        elif "ML Flow" in t or "ml flow" in t.lower():
            set_text(shape, "MLflow SSO")
        elif "натроено" in t or "DS-er" in t:
            set_text(shape, "Один аккаунт: UI + MLflow\nDS фокусируется на ML")
        elif "security check" in t.lower():
            set_text(shape, "Security pipeline")
        elif "зайти в ЛК" in t:
            set_text(shape, "Запуск pipeline из UI\nили CI на push")
        elif "HITL" in t:
            set_text(shape, "Human Approve\nдля external / tier HIGH")

    # 28 — End-to-end flow (полная 6-шаговая диаграмма)
    slide = prs.slides[27]
    for shape in slide.shapes:
        if shape.has_text_frame and "Сквозной" in shape.text_frame.text:
            set_text(shape, "Сквозной флоу")
    replace_largest_picture(slide, d / "04_sequence_flow.png", full_bleed=True)

    # 31 — CI/CD
    slide = prs.slides[30]
    for shape in slide.shapes:
        if shape.has_text_frame and "пайплайн" in shape.text_frame.text.lower():
            set_text(shape, "CI/CD pipeline")
    replace_largest_picture(slide, d / "05_cicd_pipeline.png")

    # 32 — Architecture
    slide = prs.slides[31]
    for shape in slide.shapes:
        if shape.has_text_frame and "архитектура" in shape.text_frame.text.lower():
            set_text(shape, "Архитектура")
    replace_largest_picture(slide, d / "02_compose_architecture.png")

    # 34 — Features (fix 1 layer = 1 image)
    slide = prs.slides[33]
    fixes = {
        "Детали проекта": "Особенности",
        "CI обучение": "Train в CI",
        "Неизменяемость": "Signed prod",
        "1 слой = 1 образ": "Изолированные\ngate-этапы",
        "Неподделываемая": "Attestation\nEd25519",
        "Human Approve": "Human Approve\n(исключение)",
        "lineage": "Lineage:\nданные↔код↔модель",
        "least privilege": "Sandbox для\nопасных артефактов",
        "Auth-прокси": "Keycloak +\noauth2-proxy",
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        for key, val in fixes.items():
            if key in txt:
                set_text(shape, val)
                break

    # 36 — Roadmap
    slide = prs.slides[35]
    for shape in slide.shapes:
        t = shape.text_frame.text if shape.has_text_frame else ""
        if "Направления" in t:
            set_text(shape, "Roadmap")
        elif "GenAI" in t:
            set_text(shape, "GenAI gates")
        elif "Нагрузка" in t:
            set_text(shape, "Масштабирование")
        elif "СЛОИ" in t:
            set_text(shape, "Новые gates")
        elif "SIEM" in t:
            set_text(shape, "SIEM / Jira")
        elif "Дальше" in t:
            set_text(shape, "Drift G15")

    # 37 — Analogs table
    slide = prs.slides[36]
    for shape in slide.shapes:
        if shape.has_text_frame and "Аналоги" in shape.text_frame.text:
            set_text(shape, "Аналоги — платформы vs инструменты")
    update_analogs_table(slide)
    replace_largest_picture(slide, d / "11_analogs_market.png")


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    subprocess.run([sys.executable, str(ROOT / "generate_diagrams.py")], check=True)

    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))
    apply_updates(prs)

    # Remove template junk, duplicates, sparse section slides (0-based)
    to_delete = list(range(37, len(prs.slides)))  # slides 38+
    to_delete += [
        10, 11, 14,   # lorem 11-12, duplicate benefits 15
        7,            # Use-cases section header only
        24, 25,       # sparse summary, «Идея решения»
        28, 29,       # empty admin / MLflow panels
        32, 34,       # section header «Особенности решения», «Точки роста»
    ]
    delete_slides(prs, to_delete)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print()
    print("Import to Google Slides:")
    print("  File → Import slides → Upload → FORTRESS_защита.pptx")
    print("  Or replace content in your existing doc slide-by-slide.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
